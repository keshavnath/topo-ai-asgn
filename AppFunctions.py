import pandas as pd
import json
import camelot
import fitz
from pptx import Presentation

class DataIngestion:

    @staticmethod
    # Function to read CSVs
    def read_csv(file_path):
        try:
            df = pd.read_csv(file_path)
            # df = df.interpolate(axis=1,method='linear')
            df = df.fillna(0)
            print(f"Successfully read CSV file: {file_path}")
            return df.to_dict(orient='records')
        except Exception as e:
            print(f"Error reading CSV file {file_path}: {e}")
            return {}

    @staticmethod
    # Function to read PPTXs
    def get_raw_pptx(pptx_path):
        try:
            # Load the presentation
            presentation = Presentation(pptx_path)
        except Exception as e:
            print(f"Error loading the PPTX file: {e}")
            return None, None

        # Initialize variables to store extracted data
        plaintext_data = []
        table_data = []

        # Iterate through each slide
        for i, slide in enumerate(presentation.slides):
            try:
                slide_text = []
                table_dict = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                    if shape.has_table:
                        table = shape.table
                        # Convert table to a list of lists
                        data = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text)
                            data.append(row_data)
                        # Convert to a Pandas DataFrame
                        if data:  # Ensure data is not empty
                            table = pd.DataFrame(data[1:], columns=data[0])
                            table_dict = table.to_dict(orient='records')

                        else:
                            print(f"Warning: Table on Slide {i+1} is empty.")
                        
                table_data.append((table_dict))
                plaintext_data.append((slide_text))

                print(f"Successfully read PPTX file: {pptx_path}")
                        
            
            except Exception as e:
                print(f"Error processing Slide {i + 1}: {e}")
                continue
                    
        return plaintext_data, table_data

    @staticmethod
    def parse_pptx_data(plaintext_data, table_data):
        bigdata = {}
        
        for i, slide in enumerate(plaintext_data):
            try:
                slide_dict = {}
                # slide_dict["slide_no"] = i+1
                slide_dict["title"] = slide[0]
                if table_data[i]:
                    slide_dict["table"] = table_data[i]
                if slide[1:]:
                    slide_text = {}
                    lines = slide[1].split("\n")
                    for line in lines:
                        if line[-1]==':':
                            slide_text["subtitle"] = line[:-1].strip()
                        else:

                            key, value = line.split(":", 1)
                            key = key.strip()
                            value = value.strip()
                            slide_text[key] = value
                    slide_dict["content"] = slide_text
            
                # bigdata.append(slide_dict)
                bigdata[f"slide_{i}"] = slide_dict
            
            except Exception as e:
                print(f"Error parsing slide {i+1}")
            
        # return json.dumps({"pptx_data":bigdata})
        return bigdata

    @staticmethod
    # Function to read PPTX files
    def read_pptx(file_path):
        plaintext, table = DataIngestion.get_raw_pptx(file_path)
        return DataIngestion.parse_pptx_data(plaintext, table)

    @staticmethod
    # Function to read JSON files
    def read_json(file_path):
        try:
            with open(file_path, "r") as f:
                data = json.load(f)
            print(f"Successfully read JSON file: {file_path}")
            # return json.dumps(data)
            return data

        except Exception as e:
            print(f"Error reading JSON file {file_path}: {e}")
            # return json.dumps({"json_data":[]})  # Return empty
            return None

    @staticmethod
    # Function to read PDF files
    def read_pdf(file_path):
        try:
            dfs = camelot.read_pdf(file_path, pages=str(1), flavor='lattice')
            titles = []
            tables = []
            for i in range(len(dfs)):
                try:
                    df = dfs[i].df
                    df.columns = df.iloc[0]
                    df = df[1:]
                    tables.append(df)
                    page = fitz.open(file_path).load_page(0)
                    title = page.get_text("blocks")[0][-3].strip()
                    # title = 'Table'
                    titles.append(title)
                except Exception as e:
                    print(f"Error reading Table {i+1}")
                    continue

            table_dict = {titles[i]:tables[i].to_dict(orient="records") for i in range(len(tables))}
            
            return table_dict
        
        except Exception as e:
            print(f"Error reading PDF file {file_path}: {e}")
            return None  # Return an empty on error

class DataProcessor:

    def __init__(self):
        self.file_paths = [
        "./datasets/dataset2.csv",
        "./datasets/dataset4.pptx",
        "./datasets/dataset1.json",
        "./datasets/dataset3.pdf"
        ]

        self.file_paths = dict(zip(['csv','pptx','json','pdf'],self.file_paths))

        self.file_functions = dict(zip(['csv','json','pdf','pptx'],
                                  [DataIngestion.read_csv,
                                   DataIngestion.read_json,
                                   DataIngestion.read_pdf,
                                   DataIngestion.read_pptx]))

        self.unified_data = None

    def create_unified_data(self):

        bigdata = {}

        for ft, path in self.file_paths.items():
            if self.file_functions.get(ft):
                processed = self.file_functions[ft](path)
                if processed:
                    bigdata[ft] = processed
            else:
                print(f"Could not process file of type {ft}")
        
        self.unified_data = bigdata
        # return bigdata
    
    def is_data_unified(self):
        if self.unified_data:
            return True
        else:
            return False
    

    def get_unified_data(self):
        if not (self.is_data_unified()):
            self.create_unified_data()
        
        return self.unified_data

    def get_type_data(self,ft):
        if not self.file_paths.get(ft):
                print(f"Filetype {ft} not found")
                return None
        
        else:
            if not (self.is_data_unified()):
                    return self.file_functions[ft](self.file_paths[ft])
            else:
                return self.unified_data.get(ft)